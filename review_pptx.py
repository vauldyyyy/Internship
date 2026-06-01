from pptx import Presentation

prs = Presentation(r'c:\Users\vauld\Desktop\CODING\Internship BlueStock\vauldandsouza_Submission\PPT_Slides\Final_Presentation.pptx')

print("=" * 70)
print("📊 POWERPOINT REVIEW")
print("=" * 70)
print(f"\n✅ Total Slides: {len(prs.slides)}")
print("\n📋 Slide Titles:")
print("-" * 70)

for i, slide in enumerate(prs.slides, 1):
    title = "No title"
    if slide.shapes.title:
        title = slide.shapes.title.text
    print(f"  Slide {i:2d}: {title}")

print("\n" + "=" * 70)
print("✅ PowerPoint loaded successfully!")
print("=" * 70)
